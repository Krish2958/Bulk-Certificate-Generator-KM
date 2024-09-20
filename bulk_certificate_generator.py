import os
import pandas as pd
from pptx import Presentation
import comtypes.client
from datetime import datetime
from tkinter import Tk, filedialog, messagebox, Label, Button, Frame
import tkinter as tk

def replace_placeholders(slide, replacements):
    """Replace placeholders in the slide's text with values from replacements."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    for placeholder, value in replacements.items():
                        if placeholder in run.text:
                            run.text = run.text.replace(placeholder, value)

def convert_ppt_to_pdf(input_pptx_path, output_pdf_path):
    """Convert a PPTX file to PDF format."""
    try:
        powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
        powerpoint.Visible = 1

        # Ensure the file exists before proceeding
        if not os.path.isfile(input_pptx_path):
            raise FileNotFoundError(f"The file {input_pptx_path} does not exist.")

        presentation = powerpoint.Presentations.Open(input_pptx_path, WithWindow=False)
        presentation.SaveAs(output_pdf_path, FileFormat=32)  # 32 for PDF format
        presentation.Close()
        print(f"Converted {input_pptx_path} to {output_pdf_path}")
        return True  # Indicate success
    except Exception as e:
        print(f"Failed to convert {input_pptx_path} to PDF: {e}")
        return False  # Indicate failure
    finally:
        powerpoint.Quit()

def process_ppt_and_generate_pdfs(template_ppt_path, csv_path, output_dir):
    """Generate PPTX and PDF files with placeholders replaced from CSV data."""
    # Ensure output directory exists
    os.makedirs(output_dir, exist_ok=True)

    # Read the CSV data
    data = pd.read_csv(csv_path)

    # Process each record in the CSV
    for index, row in data.iterrows():
        replacements = {f"{{{{{col}}}}}": str(row[col]) for col in data.columns}
        
        # Create a new Presentation object for each record
        prs = Presentation(template_ppt_path)
        
        # Create a copy of the presentation for each record
        timestamp = datetime.now().strftime('%Y%m%d%H%M%S')
        output_pptx_path = os.path.join(output_dir, f"{row['name']}_{timestamp}_output.pptx")
        output_pdf_path = os.path.join(output_dir, f"{row['name']}_{timestamp}_output.pdf")
        
        # Apply the replacements
        for slide in prs.slides:
            replace_placeholders(slide, replacements)

        # Save the updated PPTX and close the file
        prs.save(output_pptx_path)
        print(f"Generated PPTX: {output_pptx_path}")

        # Convert PPTX to PDF
        if convert_ppt_to_pdf(os.path.abspath(output_pptx_path), os.path.abspath(output_pdf_path)):
            # Remove the PPTX file after successful PDF conversion
            os.remove(output_pptx_path)
            print(f"Deleted PPTX: {output_pptx_path}")
        else:
            print(f"Failed to delete PPTX: {output_pptx_path} due to conversion failure.")

# UI to select files and directories
def select_ppt_template():
    global template_ppt_path
    template_ppt_path = filedialog.askopenfilename(
        title="Select PPT Template",
        filetypes=[("PowerPoint files", "*.pptx")]
    )
    template_label.config(text=f"Selected: {template_ppt_path}")

def select_csv_file():
    global csv_path
    csv_path = filedialog.askopenfilename(
        title="Select CSV File",
        filetypes=[("CSV files", "*.csv")]
    )
    csv_label.config(text=f"Selected: {csv_path}")

def select_output_dir():
    global output_dir
    output_dir = filedialog.askdirectory(title="Select Output Directory")
    output_label.config(text=f"Selected: {output_dir}")

def run_process():
    try:
        process_ppt_and_generate_pdfs(template_ppt_path, csv_path, output_dir)
        messagebox.showinfo("Success", "PPT and PDF generation completed!")
    except Exception as e:
        messagebox.showerror("Error", f"An error occurred: {e}")

# Main UI setup
root = Tk()
root.title("Bulk Certificate Generator")
root.geometry("500x400")
root.configure(bg="#f0f0f0")

# Create a frame for better layout
frame = Frame(root, bg="#ffffff", padx=20, pady=20, relief="groove", borderwidth=2)
frame.pack(padx=20, pady=20, fill="both", expand=True)

# Title label
title_label = Label(frame, text="Bulk Certificate Generator", font=("Arial", 16, "bold"), bg="#ffffff")
title_label.pack(pady=10)

# File selection labels and buttons
template_label = Label(frame, text="No PPT template selected", bg="#ffffff", fg="#333", wraplength=400)
template_label.pack(pady=5)

ppt_button = Button(frame, text="Select PPT Template", command=select_ppt_template, bg="#4CAF50", fg="white", font=("Arial", 12), width=25)
ppt_button.pack(pady=5)

csv_label = Label(frame, text="No CSV file selected", bg="#ffffff", fg="#333", wraplength=400)
csv_label.pack(pady=5)

csv_button = Button(frame, text="Select CSV File", command=select_csv_file, bg="#4CAF50", fg="white", font=("Arial", 12), width=25)
csv_button.pack(pady=5)

output_label = Label(frame, text="No output directory selected", bg="#ffffff", fg="#333", wraplength=400)
output_label.pack(pady=5)

output_button = Button(frame, text="Select Output Directory", command=select_output_dir, bg="#4CAF50", fg="white", font=("Arial", 12), width=25)
output_button.pack(pady=5)

# Generate Certificates button
run_button = Button(frame, text="Generate Certificates", command=run_process, bg="#2196F3", fg="white", font=("Arial", 14, "bold"), width=30, height=2)
run_button.pack(pady=5)

# Start the Tkinter event loop
root.mainloop()
